"""Generate slides.pptx for the QRT MDS/AML portfolio repo.

Run: python docs/_build_slides.py
Output: docs/slides.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


CHARCOAL = RGBColor(0x1A, 0x1A, 0x1F)
GRAY = RGBColor(0x55, 0x5A, 0x66)
LIGHT_GRAY = RGBColor(0xC9, 0xCD, 0xD4)
ACCENT = RGBColor(0xB8, 0x3A, 0x3A)
SUCCESS = RGBColor(0x2D, 0x6A, 0x4F)
WARN = RGBColor(0xB8, 0x80, 0x1A)
BG = RGBColor(0xFA, 0xFA, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

FIG_DIR = Path(__file__).parent / "figures"


def set_bg(slide, color=BG):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, *, size=18, color=CHARCOAL,
                bullet_color=ACCENT, spacing=Pt(6)):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = spacing
        r1 = p.add_run()
        r1.text = "■  "
        r1.font.name = FONT
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return box


def add_divider(slide, x, y, w, color=ACCENT, h_pt=2.5):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(h_pt))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_footer(slide, slide_num, total):
    add_text(
        slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
        "QRT — Overall Survival Prediction for Myeloid Leukemia · Chawki Nasrallah, 2026",
        size=10, color=GRAY,
    )
    add_text(
        slide, Inches(12.0), Inches(7.1), Inches(1.0), Inches(0.3),
        f"{slide_num} / {total}",
        size=10, color=GRAY, align=PP_ALIGN.RIGHT,
    )


def slide_title(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)

    add_divider(s, Inches(1.0), Inches(2.2), Inches(2.4), color=ACCENT, h_pt=4)
    add_text(
        s, Inches(1.0), Inches(2.5), Inches(11.5), Inches(1.6),
        "Survival Prediction under Distribution Shift",
        size=44, bold=True, color=CHARCOAL,
    )
    add_text(
        s, Inches(1.0), Inches(3.7), Inches(11.5), Inches(0.7),
        "QRT Data Challenge — Overall Survival Prediction for patients",
        size=22, color=GRAY,
    )
    add_text(
        s, Inches(1.0), Inches(4.15), Inches(11.5), Inches(0.7),
        "diagnosed with Myeloid Leukemia",
        size=22, color=GRAY,
    )
    add_text(
        s, Inches(1.0), Inches(5.7), Inches(11.5), Inches(0.4),
        "Chawki Nasrallah  ·  PhD Researcher (Biomedical Engineering, AI/Data science in Health)",
        size=14, color=CHARCOAL,
    )
    add_text(
        s, Inches(1.0), Inches(6.1), Inches(11.5), Inches(0.4),
        "github.com/chawki-nasrallah/qrt-mds-aml-survival  ·  chawki.gnasrallah@gmail.com",
        size=12, color=GRAY,
    )


def slide_tldr(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(0.5), Inches(0.5), Inches(1.5), color=ACCENT, h_pt=3)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
             "TL;DR", size=32, bold=True)

    add_text(s, Inches(0.5), Inches(1.7), Inches(12), Inches(0.5),
             "First online data-science competition.", size=18, color=GRAY, italic=True)

    add_text(s, Inches(0.5), Inches(2.4), Inches(6), Inches(0.5),
             "Public leaderboard", size=14, bold=True, color=GRAY)
    add_text(s, Inches(0.5), Inches(2.85), Inches(6), Inches(1.0),
             "0.7685", size=54, bold=True, color=SUCCESS)
    add_text(s, Inches(0.5), Inches(4.05), Inches(6), Inches(0.5),
             "Rank 11 / 743", size=20, color=CHARCOAL)

    add_text(s, Inches(7.0), Inches(2.4), Inches(6), Inches(0.5),
             "Private leaderboard", size=14, bold=True, color=GRAY)
    add_text(s, Inches(7.0), Inches(2.85), Inches(6), Inches(1.0),
             "0.7124", size=54, bold=True, color=ACCENT)
    add_text(s, Inches(7.0), Inches(4.05), Inches(6), Inches(0.5),
             "Rank 189", size=20, color=CHARCOAL)

    add_divider(s, Inches(0.5), Inches(5.5), Inches(12.3), color=LIGHT_GRAY, h_pt=1)
    add_text(s, Inches(0.5), Inches(5.7), Inches(12), Inches(1.0),
             "The gap between those two numbers is the most interesting thing in this work.",
             size=18, color=CHARCOAL, italic=True)
    add_footer(s, num, total)


def slide_problem(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(0.5), Inches(0.5), Inches(1.5), color=ACCENT, h_pt=3)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
             "The problem", size=32, bold=True)

    add_text(s, Inches(0.5), Inches(1.7), Inches(12), Inches(0.7),
             "Predict overall survival in MDS / AML patients from clinical + genetic features.",
             size=20, color=CHARCOAL)

    add_text(s, Inches(0.5), Inches(2.8), Inches(6.0), Inches(0.5),
             "Inputs", size=14, bold=True, color=GRAY)
    add_bullets(s, Inches(0.5), Inches(3.2), Inches(6.0), Inches(3.0), [
        "Clinical: BM blast %, WBC, hemoglobin, platelets",
        "Cytogenetics: ISCN-format karyotype string",
        "Molecular: per-mutation records (gene, effect, VAF, protein change)",
    ], size=15)

    add_text(s, Inches(7.0), Inches(2.8), Inches(6.0), Inches(0.5),
             "Target & metric", size=14, bold=True, color=GRAY)
    add_bullets(s, Inches(7.0), Inches(3.2), Inches(6.0), Inches(3.0), [
        "(OS_YEARS, OS_STATUS) — time-to-event survival outcome",
        "IPCW-C: censoring-aware concordance index",
        "Rewards correct pairwise risk orderings, not point predictions",
    ], size=15)

    add_divider(s, Inches(0.5), Inches(6.0), Inches(12.3), color=LIGHT_GRAY, h_pt=1)
    add_text(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.8),
             "Constraint: 2 leaderboard submissions per day.   This will matter later.",
             size=14, color=GRAY, italic=True)
    add_footer(s, num, total)


def slide_shift_features(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(0.5), Inches(0.5), Inches(1.5), color=ACCENT, h_pt=3)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
             "Distribution shift: cohort composition and features", size=28, bold=True)
    add_text(s, Inches(0.5), Inches(1.65), Inches(12), Inches(0.5),
             "Train and test were drawn from different mixtures of the same disease spectrum.",
             size=14, color=GRAY)

    s.shapes.add_picture(str(FIG_DIR / "cohort_composition.png"),
                         Inches(0.4), Inches(2.3), width=Inches(6.2))
    s.shapes.add_picture(str(FIG_DIR / "bm_blast_density.png"),
                         Inches(6.8), Inches(2.3), width=Inches(6.2))

    add_divider(s, Inches(0.5), Inches(6.3), Inches(12.3),
                color=LIGHT_GRAY, h_pt=1)
    add_text(s, Inches(0.5), Inches(6.45), Inches(12), Inches(0.5),
             "Public test has ≈4× the AML share of train (5% → 19%); median blast doubles (3% → 6%).",
             size=14, bold=True, color=CHARCOAL)
    add_text(s, Inches(0.5), Inches(6.95), Inches(12), Inches(0.4),
             "Standard covariate-shift fixes (reweighting, far-patient removal) all came back null on the LB — cross-cohort augmentation was the route that worked.",
             size=11, color=GRAY, italic=True)
    add_footer(s, num, total)


def slide_shift_outcomes(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(0.5), Inches(0.5), Inches(1.5), color=ACCENT, h_pt=3)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
             "...and so do the survival outcomes those features predict", size=26, bold=True)
    add_text(s, Inches(0.5), Inches(1.65), Inches(12), Inches(0.5),
             "AML patients die faster — so a test set enriched for AML is implicitly enriched for early-death labels.",
             size=14, color=GRAY)

    s.shapes.add_picture(str(FIG_DIR / "survival_by_cohort.png"),
                         Inches(0.5), Inches(2.3), height=Inches(3.8))

    add_text(s, Inches(7.5), Inches(2.7), Inches(5.3), Inches(0.5),
             "Train cohort (from KM curve)", size=12, bold=True, color=GRAY)
    add_bullets(s, Inches(7.5), Inches(3.2), Inches(5.3), Inches(3.0), [
        "MDS: 5-yr survival ≈ 38%",
        "AML: 5-yr survival ≈ 12%",
        "Median time-to-death (event):  MDS 1.43 y,  AML 0.86 y",
    ], size=12, spacing=Pt(10))

    add_divider(s, Inches(0.5), Inches(6.3), Inches(12.3),
                color=LIGHT_GRAY, h_pt=1)
    add_text(s, Inches(0.5), Inches(6.45), Inches(12), Inches(0.5),
             "IPCW-C rewards correct ordering among high-risk patients — the very ones the test set is enriched for.",
             size=14, bold=True, color=CHARCOAL)
    add_footer(s, num, total)


def slide_pipeline(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(0.5), Inches(0.5), Inches(1.5), color=ACCENT, h_pt=3)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
             "Pipeline: two complementary survival models", size=32, bold=True)

    box_w = Inches(5.8)
    box_h = Inches(3.6)
    box_y = Inches(1.9)

    b1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.5), box_y, box_w, box_h)
    b1.fill.solid(); b1.fill.fore_color.rgb = WHITE
    b1.line.color.rgb = LIGHT_GRAY; b1.line.width = Pt(0.75)

    b2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(7.0), box_y, box_w, box_h)
    b2.fill.solid(); b2.fill.fore_color.rgb = WHITE
    b2.line.color.rgb = LIGHT_GRAY; b2.line.width = Pt(0.75)

    add_text(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.4),
             "BASELINE  (v12)", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(0.8), Inches(2.4), Inches(5.5), Inches(0.5),
             "Wide features, single RSF", size=18, bold=True)
    add_bullets(s, Inches(0.8), Inches(3.0), Inches(5.5), Inches(2.3), [
        "~1800 features: gene one-hots × VAF, cyto one-hots, EFFECT, PROTEIN_CHANGE",
        "Dense IPSS-M score (Bernard 2022 NEJM Evidence)",
        "Random Survival Forest, 5000 trees",
    ], size=13)
    add_text(s, Inches(0.8), Inches(5.0), Inches(5.5), Inches(0.5),
             "Public LB: 0.7637", size=15, bold=True, color=SUCCESS)

    add_text(s, Inches(7.3), Inches(2.0), Inches(5.5), Inches(0.4),
             "SPECIALIST  (v22)", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(7.3), Inches(2.4), Inches(5.5), Inches(0.5),
             "23 biologically-robust features", size=18, bold=True)
    add_bullets(s, Inches(7.3), Inches(3.0), Inches(5.5), Inches(2.3), [
        "Clinical + curated cyto + curated genes + engineered scores",
        "Trained on QRT + 583 Tazi patients (propensity-selected)",
        "Single RSF (500 trees)",
    ], size=13)
    add_text(s, Inches(7.3), Inches(5.0), Inches(5.5), Inches(0.5),
             "Public LB: 0.7645 alone", size=15, bold=True, color=SUCCESS)

    add_divider(s, Inches(0.5), Inches(6.0), Inches(12.3),
                color=LIGHT_GRAY, h_pt=1)
    add_text(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.6),
             "FINAL  =  0.30 · rank(baseline)  +  0.70 · rank(specialist)    →  Public LB: 0.7685",
             size=16, bold=True, color=CHARCOAL)
    add_text(s, Inches(0.5), Inches(6.75), Inches(12), Inches(0.4),
             "rank-blend (not raw): two RSF architectures live on different scales, ~[218,2039] vs ~[127,1204]",
             size=11, color=GRAY, italic=True)
    add_footer(s, num, total)


def slide_baseline_detail(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(0.5), Inches(0.5), Inches(1.5), color=ACCENT, h_pt=3)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
             "Baseline detail: feature engineering", size=32, bold=True)
    add_text(s, Inches(0.5), Inches(1.55), Inches(12), Inches(0.5),
             "Combining five sources of signal into ~1800 columns.",
             size=14, color=GRAY)

    add_bullets(s, Inches(0.5), Inches(2.2), Inches(12.5), Inches(5.0), [
        "Clinical numerics — log1p transform compresses right-skewed counts (WBC, PLT) so the tree splits aren't dominated by extreme values;  MICE imputation iteratively regresses each variable on the others (RandomForest as regressor) to preserve cross-feature correlations that median-fill would destroy.",
        "Per-gene one-hot weighted by max VAF — captures both presence (was the gene mutated?) and allelic burden (how clonal is the mutation?).",
        "Two cytogenetic parsers — a broad regex extracts every +/−, translocation and inversion (high recall, surfaces rare lesions); a curated parser produces clean IPSS-R-aligned flags (high precision).",
        "EFFECT and PROTEIN_CHANGE — mutation type (missense, ITD, PTD, …) and residue position parsed from HGVS strings like p.R882H.",
        "IPSS-M score (Bernard 2022) — the full Molecular International Prognostic Scoring System as a single dense feature.  Carries ≈42% of tree-ensemble importance.",
        "Two interaction flags — NPM1+/FLT3− (favorable AML), and TP53+ × complex karyotype (high-risk MDS/AML).",
    ], size=12, spacing=Pt(7))
    add_footer(s, num, total)


def slide_specialist_detail(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(0.5), Inches(0.5), Inches(1.5), color=ACCENT, h_pt=3)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
             "Specialist: cross-cohort augmentation via propensity selection", size=26, bold=True)
    add_text(s, Inches(0.5), Inches(1.6), Inches(12), Inches(0.5),
             "Pull external patients matched to the QRT high-blast test region.",
             size=14, color=GRAY)

    add_text(s, Inches(0.5), Inches(2.3), Inches(12), Inches(0.5),
             "How the propensity classifier works", size=13, bold=True, color=GRAY)
    add_bullets(s, Inches(0.5), Inches(2.7), Inches(12.5), Inches(1.9), [
        "Label QRT test patients with BM blast ≥ 20% as 1; label all Tazi 2022 NEJM Evidence cohort patients as 0",
        "Train LightGBM 5-fold to discriminate, on raw inputs only (no engineered scores)",
        "Sort Tazi by P(label=1); keep the top 583 — most similar to underserved AML region",
    ], size=12, spacing=Pt(6))

    add_text(s, Inches(0.5), Inches(4.65), Inches(6.0), Inches(0.4),
             "Why K = 583", size=13, bold=True, color=ACCENT)
    add_text(s, Inches(0.5), Inches(5.05), Inches(6.0), Inches(2.0),
             "Sized to match the public test's AML share.  Adding 583 high-blast Tazi patients to the 3,173 QRT train cohort lifts the augmented train AML share from ≈5% to ≈19% — the same AML share as the public test set.",
             size=11, color=CHARCOAL)

    add_text(s, Inches(7.0), Inches(4.65), Inches(6.0), Inches(0.4),
             "Why raw inputs only", size=13, bold=True, color=ACCENT)
    add_text(s, Inches(7.0), Inches(5.05), Inches(6.0), Inches(2.0),
             "Engineered features (IPSS-M, TP53 multihit) inside the propensity model would correlate selection criteria with the specialist's own scoring — a feedback loop that narrows the augmented distribution. Worth +0.0011 public LB.",
             size=11, color=CHARCOAL)
    add_footer(s, num, total)


def slide_blend(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(0.5), Inches(0.5), Inches(1.5), color=ACCENT, h_pt=3)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
             "Rank-blending: scale-invariant complementarity", size=32, bold=True)

    add_text(s, Inches(0.5), Inches(1.8), Inches(12), Inches(0.5),
             "The two models output incompatible scales:", size=16, color=CHARCOAL)

    add_text(s, Inches(1.5), Inches(2.6), Inches(5), Inches(0.5),
             "Baseline range:", size=14, color=GRAY)
    add_text(s, Inches(4.5), Inches(2.55), Inches(5), Inches(0.5),
             "~[218,  2039]", size=18, bold=True, color=CHARCOAL, font="Consolas")
    add_text(s, Inches(1.5), Inches(3.2), Inches(5), Inches(0.5),
             "Specialist range:", size=14, color=GRAY)
    add_text(s, Inches(4.5), Inches(3.15), Inches(5), Inches(0.5),
             "~[127,  1204]", size=18, bold=True, color=CHARCOAL, font="Consolas")

    add_text(s, Inches(0.5), Inches(4.0), Inches(12), Inches(0.7),
             "→ Raw averaging lets the baseline dominate by scale alone.",
             size=16, color=ACCENT)

    add_divider(s, Inches(0.5), Inches(4.9), Inches(12.3),
                color=LIGHT_GRAY, h_pt=1)
    add_text(s, Inches(0.5), Inches(5.1), Inches(12), Inches(0.7),
             "Convert to ranks, then blend.",
             size=18, bold=True, color=CHARCOAL)
    add_text(s, Inches(0.5), Inches(5.7), Inches(12), Inches(0.6),
             "final  =  0.30 · rank(baseline)  +  0.70 · rank(specialist)",
             size=22, bold=True, color=SUCCESS, font="Consolas")
    add_text(s, Inches(0.5), Inches(6.4), Inches(12), Inches(0.5),
             "Also matches IPCW-C's geometry — the metric scores pairwise orderings, not magnitudes.",
             size=12, color=GRAY, italic=True)
    add_footer(s, num, total)


def slide_ablation(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(0.5), Inches(0.5), Inches(1.5), color=ACCENT, h_pt=3)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
             "Ablation: the augmentation is the load-bearing element", size=28, bold=True)
    add_text(s, Inches(0.5), Inches(1.65), Inches(12), Inches(0.5),
             "Each piece earning its keep on the public LB.",
             size=14, color=GRAY, italic=True)

    headers = ["Variant", "Spec training", "α", "Public LB", "Δ vs v12"]
    rows = [
        ("v12 baseline alone",          "—",                "—",    "0.7637", "—",      CHARCOAL, False),
        ("Specialist alone",             "QRT + Tazi-583",   "—",    "0.7645", "+0.0008", CHARCOAL, False),
        ("Spec + v12, NO augmentation",  "QRT only",         "0.50", "0.7584", "−0.0053", ACCENT,   True),
        ("Spec + v12, with augmentation","QRT + Tazi-583",   "0.50", "0.7674", "+0.0037", CHARCOAL, False),
        ("Final (v22, Option C)",        "QRT + Tazi-583",   "0.30", "0.7685", "+0.0048", SUCCESS,  True),
    ]
    table_y = Inches(2.3)
    col_w = [Inches(3.5), Inches(2.8), Inches(0.9), Inches(1.5), Inches(1.5)]
    x_start = Inches(0.6)

    x = x_start
    for i, h in enumerate(headers):
        align = PP_ALIGN.RIGHT if i >= 2 else PP_ALIGN.LEFT
        add_text(s, x, table_y, col_w[i], Inches(0.5), h,
                 size=12, bold=True, color=GRAY, align=align)
        x += col_w[i]
    add_divider(s, x_start, table_y + Inches(0.55),
                Inches(10.2), color=LIGHT_GRAY, h_pt=1)

    for r, (name, train, alpha, lb, delta, color, bold) in enumerate(rows):
        row_y = table_y + Inches(0.75) + Inches(0.5) * r
        cells = [name, train, alpha, lb, delta]
        x = x_start
        for i, cell in enumerate(cells):
            font = "Consolas" if i >= 2 else FONT
            align = PP_ALIGN.RIGHT if i >= 2 else PP_ALIGN.LEFT
            cell_color = color if i in (3, 4) else CHARCOAL
            cell_bold = bold if i in (0, 3, 4) else False
            add_text(s, x, row_y, col_w[i], Inches(0.5), cell,
                     size=12, color=cell_color, bold=cell_bold, font=font, align=align)
            x += col_w[i]

    add_divider(s, Inches(0.5), Inches(5.8), Inches(12.3),
                color=LIGHT_GRAY, h_pt=1)
    add_text(s, Inches(0.5), Inches(5.95), Inches(12), Inches(0.5),
             "Apples-to-apples at α=0.50:  Tazi augmentation alone moves the blend +0.0090 (0.7584 → 0.7674).",
             size=14, bold=True, color=CHARCOAL)
    add_text(s, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
             "Without augmentation, the specialist drags the blend BELOW v12 alone (−0.0053).",
             size=12, color=ACCENT, italic=True)
    add_footer(s, num, total)


def slide_results_public(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(0.5), Inches(0.5), Inches(1.5), color=ACCENT, h_pt=3)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
             "Public leaderboard: rank 11 / 743", size=28, bold=True)

    headers = ["Rank", "Score", "Team"]
    rows = [
        ("1",  "0.7739", "ocercle"),
        ("2",  "0.7732", "qrtmieuxquecitadel"),
        ("3",  "0.7732", "Astra"),
        ("…",  "…",      "…"),
        ("11", "0.7685", "chawkinas  (me)"),
    ]
    table_y = Inches(2.2)
    col_w = [Inches(1.4), Inches(2.2), Inches(8.0)]
    x_start = Inches(1.0)

    x = x_start
    for i, h in enumerate(headers):
        add_text(s, x, table_y, col_w[i], Inches(0.5), h,
                 size=14, bold=True, color=GRAY)
        x += col_w[i]
    add_divider(s, x_start, table_y + Inches(0.55),
                Inches(11.6), color=LIGHT_GRAY, h_pt=1)

    for r, row in enumerate(rows):
        row_y = table_y + Inches(0.75) + Inches(0.6) * r
        x = x_start
        is_me = r == len(rows) - 1
        for i, cell in enumerate(row):
            color = SUCCESS if is_me else CHARCOAL
            bold = is_me
            font = "Consolas" if i in (0, 1) else FONT
            add_text(s, x, row_y, col_w[i], Inches(0.5), cell,
                     size=16, color=color, bold=bold, font=font)
            x += col_w[i]

    add_divider(s, Inches(0.5), Inches(6.4), Inches(12.3),
                color=LIGHT_GRAY, h_pt=1)
    add_text(s, Inches(0.5), Inches(6.55), Inches(12), Inches(0.5),
             "Top of the public leaderboard at 0.7739 — within 0.0054 of my score.",
             size=12, color=GRAY, italic=True)
    add_footer(s, num, total)


def slide_results_private(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(0.5), Inches(0.5), Inches(1.5), color=ACCENT, h_pt=3)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
             "Private leaderboard: a 5-point shake-up across the field", size=28, bold=True)

    headers = ["", "Public", "Private", "Drop"]
    rows = [
        ("Top of LB (rank 1)",     "0.7739", "0.7231", "−0.0508"),
        ("chawkinas  (me)",        "0.7685", "0.7124", "−0.0561"),
    ]
    table_y = Inches(2.3)
    col_w = [Inches(4.5), Inches(2.0), Inches(2.0), Inches(2.0)]
    x_start = Inches(1.0)

    x = x_start
    for i, h in enumerate(headers):
        add_text(s, x, table_y, col_w[i], Inches(0.5), h,
                 size=14, bold=True, color=GRAY)
        x += col_w[i]
    add_divider(s, x_start, table_y + Inches(0.55),
                Inches(10.5), color=LIGHT_GRAY, h_pt=1)

    for r, row in enumerate(rows):
        row_y = table_y + Inches(0.85) + Inches(0.75) * r
        x = x_start
        is_me = r == 1
        for i, cell in enumerate(row):
            font = "Consolas" if i in (1, 2, 3) else FONT
            color = ACCENT if (i == 3) else (SUCCESS if is_me and i == 0 else CHARCOAL)
            bold = is_me or i == 3
            add_text(s, x, row_y, col_w[i], Inches(0.5), cell,
                     size=16, color=color, bold=bold, font=font)
            x += col_w[i]

    add_divider(s, Inches(0.5), Inches(5.5), Inches(12.3),
                color=LIGHT_GRAY, h_pt=1)
    add_text(s, Inches(0.5), Inches(5.65), Inches(12), Inches(0.5),
             "The entire field shook up — top public 0.7739 → top private 0.7231 (−0.0508).",
             size=14, color=CHARCOAL)
    add_text(s, Inches(0.5), Inches(6.15), Inches(12), Inches(0.5),
             "My drop is 0.0053 worse than the winner's — within ~0.005 of the rest of the top of the public LB.",
             size=12, color=GRAY, italic=True)
    add_text(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
             "For reference: organisers' benchmark scored 0.6411 on private (rank 678) — beat it by +0.0713.",
             size=12, color=GRAY)
    add_footer(s, num, total)


def slide_early_submitters(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(0.5), Inches(0.5), Inches(1.5), color=ACCENT, h_pt=3)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
             "Private-LB winners stopped iterating early", size=32, bold=True)
    add_text(s, Inches(0.5), Inches(1.65), Inches(12), Inches(0.5),
             "Top-5 private finishers submitted months before the public LB became dense to mine.",
             size=14, color=GRAY)

    headers = ["Private rank", "Score", "Submission date", "Team"]
    rows = [
        ("1",  "0.7231", "Dec 14, 2025",  "arthur_derouck"),
        ("2",  "0.7231", "Dec 14, 2025",  "arthur_derouck & rbarata"),
        ("3",  "0.7219", "Jan 4, 2026",   "cocohhhhh"),
        ("4",  "0.7216", "Aug 15, 2025",  "djtiesto"),
        ("5",  "0.7208", "Mar 1, 2025",   "guppsFTSF"),
    ]
    table_y = Inches(2.5)
    col_w = [Inches(1.8), Inches(1.8), Inches(2.8), Inches(5.5)]
    x_start = Inches(0.7)

    x = x_start
    for i, h in enumerate(headers):
        add_text(s, x, table_y, col_w[i], Inches(0.5), h,
                 size=13, bold=True, color=GRAY)
        x += col_w[i]
    add_divider(s, x_start, table_y + Inches(0.55),
                Inches(11.9), color=LIGHT_GRAY, h_pt=1)

    for r, row in enumerate(rows):
        row_y = table_y + Inches(0.75) + Inches(0.5) * r
        x = x_start
        for i, cell in enumerate(row):
            font = "Consolas" if i in (0, 1) else FONT
            color = ACCENT if i == 2 and ("Mar 2025" in cell or "Aug 15" in cell) else CHARCOAL
            bold = i == 2 and ("Mar 2025" in cell or "Aug 15" in cell)
            add_text(s, x, row_y, col_w[i], Inches(0.5), cell,
                     size=13, color=color, bold=bold, font=font)
            x += col_w[i]

    add_divider(s, Inches(0.5), Inches(6.0), Inches(12.3),
                color=LIGHT_GRAY, h_pt=1)
    add_text(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.6),
             "Late-June 2026 submitters (including me) uniformly suffered the BIGGEST private drops.",
             size=14, color=ACCENT, bold=True)
    add_text(s, Inches(0.5), Inches(6.65), Inches(12), Inches(0.5),
             "Goodhart's law on the public test set: optimize against it long enough and you encode its noise.",
             size=12, color=GRAY, italic=True)
    add_footer(s, num, total)


def slide_tech_lessons(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(0.5), Inches(0.5), Inches(1.5), color=ACCENT, h_pt=3)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
             "Technical lessons", size=32, bold=True)
    add_text(s, Inches(0.5), Inches(1.55), Inches(12), Inches(0.5),
             "Three findings that transfer beyond this project.",
             size=14, color=GRAY, italic=True)

    sections = [
        (
            "1.  Components that fail individually can win in combination.",
            "Specialist alone −0.0040 LB · Tazi augmentation alone −0.0033 · combined →  +0.0048 over v12.",
            "Don't drop a building block just because it loses standalone — multi-component systems with complementary errors can beat the sum of their parts.",
        ),
        (
            "2.  Feature-schema design enables cross-cohort augmentation.",
            "Same 583 Tazi patients: into baseline's 1,800 sparse features →  −0.0033 LB;  into specialist's 23 shared features →  +0.0090 swing.",
            "More data isn't enough — the feature vocabulary has to transfer across cohorts. Plan for transfer at the feature stage, not the model stage.",
        ),
        (
            "3.  Beware feedback loops in propensity-driven selection.",
            "v21 propensity used engineered scores (the specialist's own signals);  v22 used raw inputs only;  +0.0011 LB from breaking the loop.",
            "If a selection procedure shares features with the model it serves, the model is implicitly choosing its own training data — a subtle selection bias.",
        ),
    ]

    y0 = 2.2
    section_h = 1.65
    for i, (head, evidence, principle) in enumerate(sections):
        y = y0 + i * section_h
        add_text(s, Inches(0.5), Inches(y), Inches(12.5), Inches(0.4),
                 head, size=15, bold=True, color=CHARCOAL)
        add_text(s, Inches(0.7), Inches(y + 0.5), Inches(12.3), Inches(0.4),
                 evidence, size=11, color=GRAY)
        add_text(s, Inches(0.7), Inches(y + 0.95), Inches(12.3), Inches(0.5),
                 principle, size=12, color=CHARCOAL, italic=True)

    add_footer(s, num, total)


def slide_lessons(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(0.5), Inches(0.5), Inches(1.5), color=ACCENT, h_pt=3)
    add_text(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
             "What I learned", size=32, bold=True)
    add_text(s, Inches(0.5), Inches(1.6), Inches(12), Inches(0.5),
             "First DS competition.  The methodology I'd defend; the iteration policy I'd change.",
             size=14, color=GRAY, italic=True)

    add_bullets(s, Inches(0.5), Inches(2.5), Inches(12), Inches(4.5), [
        "Late-stage optimization against a public test set is itself a form of overfitting.",
        "When CV and LB disagree, trust CV — especially once per-iteration LB deltas drop below noise band (~0.003).",
        "Pick a final submission once deltas plateau; don't nudge it for marginal LB gains.",
        "Treat the public LB as a sanity check, not a continuous optimization target.",
        "A 0.001 improvement on a 1193-patient C-index is within sampling noise.  The improvement may not be real.",
    ], size=16, spacing=Pt(12))
    add_footer(s, num, total)


def slide_close(prs, total, num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    add_divider(s, Inches(1.0), Inches(2.6), Inches(2.4), color=ACCENT, h_pt=4)
    add_text(s, Inches(1.0), Inches(2.9), Inches(11.5), Inches(1.2),
             "Thank you.", size=44, bold=True, color=CHARCOAL)
    add_text(s, Inches(1.0), Inches(4.1), Inches(11.5), Inches(0.5),
             "Code, methodology writeup, post-mortem:", size=16, color=GRAY)
    add_text(s, Inches(1.0), Inches(4.6), Inches(11.5), Inches(0.5),
             "github.com/chawki-nasrallah/qrt-mds-aml-survival", size=18, bold=True, color=CHARCOAL, font="Consolas")
    add_text(s, Inches(1.0), Inches(5.6), Inches(11.5), Inches(0.5),
             "Chawki Nasrallah  ·  chawki.gnasrallah@gmail.com", size=14, color=CHARCOAL)


def build(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_title, slide_tldr, slide_problem,
        slide_shift_features, slide_shift_outcomes,
        slide_pipeline, slide_baseline_detail, slide_specialist_detail, slide_blend,
        slide_ablation,
        slide_results_public, slide_results_private,
        slide_tech_lessons, slide_lessons,
        slide_close,
    ]
    total = len(builders)

    slide_title(prs, total)
    for i, builder in enumerate(builders[1:-1], start=2):
        builder(prs, total, i)
    slide_close(prs, total, total)

    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    build(Path(__file__).parent / "slides.pptx")
